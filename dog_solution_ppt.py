from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# 발표자/날짜 입력 (수정 가능)
presenter_name = "한재원"
today = datetime.now().strftime("%Y-%m-%d")

slides_data = [
    # 1. 제목 슬라이드
    {
        "title": "강아지 만능 솔루션 🐾",
        "content": f"반려견을 위한 AI 기반 통합 관리 플랫폼\n\n발표자: {presenter_name}\n날짜: {today}"
    },
    # 2. 앱 소개 및 개요
    {
        "title": "\"강아지 만능 솔루션\" 이란?",
        "content": (
            "- 반려견 보호자를 위한 올인원(All-in-one) 웹 애플리케이션\n"
            "- AI 기술을 활용하여 반려견의 품종 분석부터 건강 관리, 정보 탐색까지 다양한 기능 제공\n"
            "- Streamlit 프레임워크를 사용하여 쉽고 빠르게 개발 및 배포 가능"
        )
    },
    # 3. 견종 분석기
    {
        "title": "🔍 우리 강아지는 무슨 품종일까? - 견종 분석기",
        "content": (
            "- 강아지 이미지 업로드 시, AI 모델(ResNet50 기반)이 품종을 예측하고 상위 3개 품종의 확률 제시\n"
            "- 유기견 품종 파악, 믹스견 품종 추정 등\n"
            "- [이미지 업로드/예측 결과 화면 스크린샷]"
        )
    },
    # 4. 전문가 상담
    {
        "title": "💬 궁금증 해결! - 전문가 상담 (AI 챗봇)",
        "content": (
            "- OpenAI API 기반 AI 챗봇, 반려견 관련 전문 답변 제공\n"
            "- 사료, 훈련, 건강 등 다양한 주제\n"
            "- [챗봇 대화 화면 스크린샷]"
        )
    },
    # 5. 주변 장소 찾기
    {
        "title": "🗺️ 반려견과 함께! - 주변 장소 찾기",
        "content": (
            "- Google Maps API로 애견동반 카페/식당/동물병원/공원/숙소 등 검색 및 지도 표시\n"
            "- 외출 시 유용한 장소 정보 제공\n"
            "- API 키 문제 해결로 안정적 서비스\n"
            "- [지도 및 검색결과 화면 스크린샷]"
        )
    },
    # 6. 성장/건강 일지
    {
        "title": "📔 우리 아이 성장 기록 - 성장/건강 일지",
        "content": (
            "- 날짜별 체중, 식사량, 특이사항 기록 및 관리\n"
            "- 건강 변화 추이 파악, 병원 방문 기초자료\n"
            "- [일지 입력/기록, 그래프 화면 스크린샷]"
        )
    },
    # 7. 기타 기능들
    {
        "title": "🎁🎨❓🗣️ 즐거움을 더하는 기능들",
        "content": (
            "- 품종 비교, 맞춤 추천\n"
            "- AI 이미지 생성 (예시 이미지)\n"
            "- 강아지 퀴즈 게임\n"
            "- 행동/소리 분석 (체험형)\n"
            "- [각 기능 스크린샷]"
        )
    },
    # 8. 기술 스택
    {
        "title": "🛠️ 기술 스택",
        "content": (
            "- Streamlit (UI 개발/배포)\n"
            "- PyTorch(ResNet50), OpenAI API(GPT)\n"
            "- pandas\n"
            "- Google Maps Platform API (Geocoding, Places)\n"
            "- PIL, torchvision, speech_recognition, gTTS 등"
        )
    },
    # 9. 향후 계획
    {
        "title": "🚀 향후 계획",
        "content": (
            "- 사용자 피드백 기반 기능 개선\n"
            "- 모바일 최적화 및 반응형 디자인\n"
            "- 행동 인식·질병 진단 등 AI 모델 확대\n"
            "- 커뮤니티 기능(정보 공유)"
        )
    },
    # 10. Q&A/감사
    {
        "title": "질문 및 답변",
        "content": (
            "경청해 주셔서 감사합니다.\n\n질문 있으시면 해주세요.\n[연락처/이메일]"
        )
    },
]

prs = Presentation()
for idx, slide in enumerate(slides_data):
    layout = prs.slide_layouts[0] if idx == 0 else prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    title = s.shapes.title
    title.text = slide["title"]
    content_shape = s.placeholders[1]
    content_shape.text = slide["content"]
    for paragraph in content_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20 if idx == 0 else 18)
            run.font.name = "맑은 고딕"

prs.save("강아지_만능_솔루션_발표.pptx")
print("파워포인트 파일이 생성되었습니다! (강아지_만능_솔루션_발표.pptx)")
